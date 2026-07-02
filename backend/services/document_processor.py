# backend/services/document_processor.py

import io
from pypdf import PdfReader
import docx
from pptx import Presentation
from PIL import Image
import pytesseract

def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract text from PDF file bytes using pypdf."""
    pdf_file = io.BytesIO(file_bytes)
    reader = PdfReader(pdf_file)
    text = ""
    for page in reader.pages:
        extracted = page.extract_text()
        if extracted:
            text += extracted + "\n"
    return text

def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract text from DOCX file bytes using python-docx."""
    docx_file = io.BytesIO(file_bytes)
    doc = docx.Document(docx_file)
    text = ""
    for para in doc.paragraphs:
        if para.text:
            text += para.text + "\n"
    return text

def extract_text_from_txt(file_bytes: bytes) -> str:
    """Extract text from TXT file bytes supporting UTF-8 and Latin-1."""
    try:
        return file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        return file_bytes.decode("latin-1")

def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract text and speaker notes from PowerPoint slideshow using python-pptx."""
    prs = Presentation(io.BytesIO(file_bytes))
    text = ""
    for slide_idx, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)
        
        # Include speaker notes if available
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
            if notes:
                slide_text.append(f"[Speaker Notes]: {notes}")
                
        if slide_text:
            text += f"--- Slide {slide_idx + 1} ---\n" + "\n".join(slide_text) + "\n\n"
    return text

def extract_text_from_image(file_bytes: bytes) -> str:
    """Extract text from image bytes (PNG, JPEG, WEBP, TIFF) using pytesseract OCR."""
    try:
        image = Image.open(io.BytesIO(file_bytes))
        text = pytesseract.image_to_string(image)
        return text
    except Exception as e:
        # Graceful handling for missing local tesseract binary
        raise RuntimeError(
            f"OCR processing failed. Ensure Tesseract OCR engine is installed on the host system. Error: {str(e)}"
        )

def extract_text(file_bytes: bytes, filename: str) -> str:
    """Unified function to extract text based on file extension."""
    ext = filename.split(".")[-1].lower()
    if ext == "pdf":
        return extract_text_from_pdf(file_bytes)
    elif ext == "docx":
        return extract_text_from_docx(file_bytes)
    elif ext in ["txt", "text"]:
        return extract_text_from_txt(file_bytes)
    elif ext in ["pptx", "ppt"]:
        return extract_text_from_pptx(file_bytes)
    elif ext in ["png", "jpg", "jpeg", "webp", "tiff", "tif"]:
        return extract_text_from_image(file_bytes)
    else:
        raise ValueError(f"Unsupported file format: .{ext}")
